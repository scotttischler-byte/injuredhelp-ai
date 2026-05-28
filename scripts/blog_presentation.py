#!/usr/bin/env python3
"""Generate gold-tier PowerPoint summaries for WreckMatch blog posts (SEO / GEO / indexing)."""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
BLOG_DIR = ROOT / "content/blog"
BLOG_ES_DIR = ROOT / "content/blog/es"
PRESENTATIONS_DIR = ROOT / "public/blog/presentations"

COPY: dict[str, dict[str, str]] = {
    "en": {
        "read_first": "Important — read first",
        "quick_answer": "Quick answer (quotable summary)",
        "faq_heading": "Frequently asked questions",
        "cta_heading": "Get matched with a licensed attorney",
        "review_heading": "Reviewed for legal context",
        "disclaimer_1": "Educational only — not legal advice.",
        "disclaimer_2": "WreckMatch LLC is a legal referral service, NOT a law firm.",
        "disclaimer_3": "800+ participating law firms nationwide — free matching.",
        "disclaimer_4": "Confirm all deadlines with licensed counsel in your state.",
        "cta_1": "Free attorney matching in ~60 seconds: {site}/#form",
        "cta_2": "Call {phone}",
        "cta_3": "Full article: {url}",
        "cta_4": "No fee for matching — contingency attorneys in network.",
    },
    "es": {
        "read_first": "Importante — lea primero",
        "quick_answer": "Respuesta rápida (resumen citabile)",
        "faq_heading": "Preguntas frecuentes",
        "cta_heading": "Emparejamiento con abogado con licencia",
        "review_heading": "Revisado para contexto legal",
        "disclaimer_1": "Solo educativo — no es asesoría legal.",
        "disclaimer_2": "WreckMatch LLC es un servicio de referencia legal, NO un bufete.",
        "disclaimer_3": "Más de 800 bufetes participantes — emparejamiento gratuito.",
        "disclaimer_4": "Confirme plazos con un abogado con licencia en su estado.",
        "cta_1": "Emparejamiento gratis en ~60 segundos: {site}/#form",
        "cta_2": "Llame al {phone}",
        "cta_3": "Artículo completo: {url}",
        "cta_4": "Sin costo por emparejamiento — honorarios de contingencia.",
    },
}
SITE = "https://www.wreckmatch.com"
PHONE = "855 WRECKMATCH (855) 897-3256"
NOTES_FOOTER = (
    " WreckMatch LLC — legal referral service, not a law firm. Educational only. "
    "800+ participating law firms. Call 855 WRECKMATCH. "
    "Reviewed for legal context by Judge Roy Waddell."
)


def _notes(text: str) -> str:
    return (text + NOTES_FOOTER)[:1500]

WM_RED = RGBColor(0xCC, 0x00, 0x00)
WM_DARK = RGBColor(0x0F, 0x17, 0x2A)
WM_EMERALD = RGBColor(0x10, 0xB9, 0x81)
WM_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WM_SLATE = RGBColor(0x94, 0xA3, 0xB8)

SKIP_HEADINGS = re.compile(
    r"(related resources|frequently asked questions|accident survival guide|trust, compliance)",
    re.I,
)


@dataclass
class SlideContent:
    heading: str
    bullets: list[str]
    notes: str = ""


@dataclass
class PresentationReport:
    slug: str
    score: int
    slide_count: int
    issues: list[str] = field(default_factory=list)
    path: Path | None = None

    @property
    def ok(self) -> bool:
        return self.score >= 100


def parse_frontmatter(text: str) -> tuple[dict[str, str], str]:
    if not text.startswith("---"):
        return {}, text
    parts = text.split("---", 2)
    if len(parts) < 3:
        return {}, text
    raw, body = parts[1], parts[2]
    fm: dict[str, str] = {}

    def folded(key: str) -> str | None:
        m = re.search(rf"^{key}:\s*>\-?\s*\n((?:  .+\n?)+)", raw, re.M)
        if not m:
            return None
        return " ".join(line.strip() for line in m.group(1).splitlines() if line.strip())

    for key in ("title", "description", "excerpt"):
        v = folded(key)
        if v:
            fm[key] = v
    for line in raw.strip().splitlines():
        if ":" not in line or line.startswith(" "):
            continue
        k, v = line.split(":", 1)
        k = k.strip()
        if k in fm:
            continue
        v = v.strip().strip('"').strip("'")
        if v in (">-", ">", "|"):
            continue
        fm[k] = v
    return fm, body


def strip_md(text: str) -> str:
    text = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"<!--[\s\S]*?-->", "", text)
    return text.strip()


def extract_sections(body: str) -> list[SlideContent]:
    body = re.split(r"## Related resources", body, 1)[0]
    sections: list[SlideContent] = []
    chunks = re.split(r"\n(?=## )", body)
    for chunk in chunks:
        chunk = chunk.strip()
        if not chunk.startswith("## "):
            continue
        lines = chunk.split("\n")
        heading = strip_md(lines[0].lstrip("# ").strip())
        if SKIP_HEADINGS.search(heading):
            continue
        rest = "\n".join(lines[1:]).strip()
        bullets: list[str] = []
        notes_parts: list[str] = []
        if "### " in rest:
            faq_blocks = re.split(r"\n(?=### )", rest)
            for block in faq_blocks:
                block = block.strip()
                if block.startswith("### "):
                    q = strip_md(block.split("\n", 1)[0].lstrip("# ").strip())
                    a = strip_md(block.split("\n", 1)[1] if "\n" in block else "")
                    if q and a:
                        bullets.append(f"Q: {q}")
                        bullets.append(f"A: {a[:280]}")
                        notes_parts.append(f"{q} — {a}")
                else:
                    bullets.extend(_bullets_from_block(block))
        else:
            bullets.extend(_bullets_from_block(rest))
            notes_parts.append(strip_md(rest)[:1200])
        if not bullets:
            para = strip_md(rest)[:400]
            if para:
                bullets = [para]
        bullets = [b[:320] for b in bullets if b.strip()][:8]
        if bullets:
            sections.append(
                SlideContent(
                    heading=heading[:120],
                    bullets=bullets,
                    notes=" ".join(notes_parts)[:1500] or heading,
                )
            )
    return sections


def _bullets_from_block(block: str) -> list[str]:
    out: list[str] = []
    for line in block.splitlines():
        line = line.strip()
        if re.match(r"^[-*]\s+", line) or re.match(r"^\d+\.\s+", line):
            out.append(strip_md(re.sub(r"^[-*]\s+|^\d+\.\s+", "", line)))
        elif line.startswith("|") and "---" not in line:
            cells = [strip_md(c) for c in line.split("|") if c.strip()]
            if cells:
                out.append(" · ".join(cells)[:300])
    if not out:
        para = strip_md(block)
        if para and len(para) > 40:
            out.append(para[:400])
    return out


def find_quick_answer(body: str) -> str:
    m = re.search(
        r"\*\*(?:Quick answer|En resumen|Respuesta rápida):\*\*\s*(.+?)(?:\n\n|\n## )",
        body,
        re.S | re.I,
    )
    if m:
        return strip_md(m.group(1))[:500]
    m = re.search(r"\*\*Direct answer:\*\*\s*(.+?)(?:\n\n|\n## )", body, re.S | re.I)
    if m:
        return strip_md(m.group(1))[:500]
    m = re.search(r"\*\*At a glance:\*\*\s*(.+?)(?:\n\n|\n## )", body, re.S | re.I)
    if m:
        return strip_md(m.group(1))[:500]
    return ""


def faq_pairs(body: str) -> list[tuple[str, str]]:
    pairs: list[tuple[str, str]] = []
    faq_zone = body
    for marker in ("## Frequently asked questions", "## Preguntas frecuentes"):
        if marker in body:
            faq_zone = body.split(marker, 1)[1]
            break
    for m in re.finditer(r"###\s+(.+?)\n+([\s\S]*?)(?=\n### |\n## |\Z)", faq_zone):
        q = strip_md(m.group(1).strip())
        a = strip_md(m.group(2).strip())[:400]
        if q and a and len(pairs) < 6:
            pairs.append((q, a))
    return pairs


def _set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_title_slide(prs: Presentation, title: str, subtitle: str, notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WM_DARK)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8.8), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "WreckMatch.com"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WM_EMERALD
    p.alignment = PP_ALIGN.LEFT

    box2 = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(8.8), Inches(2.4))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title[:200]
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = WM_WHITE

    if subtitle:
        p3 = tf2.add_paragraph()
        p3.text = subtitle[:300]
        p3.font.size = Pt(16)
        p3.font.color.rgb = WM_SLATE
        p3.space_before = Pt(12)

    slide.notes_slide.notes_text_frame.text = _notes(notes)


def _add_content_slide(prs: Presentation, heading: str, bullets: list[str], notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WM_DARK)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WM_RED
    bar.line.fill.background()

    hbox = slide.shapes.add_textbox(Inches(0.45), Inches(0.45), Inches(9.0), Inches(0.9))
    hp = hbox.text_frame.paragraphs[0]
    hp.text = heading
    hp.font.size = Pt(24)
    hp.font.bold = True
    hp.font.color.rgb = WM_WHITE

    bbox = slide.shapes.add_textbox(Inches(0.55), Inches(1.35), Inches(8.9), Inches(5.5))
    tf = bbox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets[:7]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = bullet
        para.level = 0
        para.font.size = Pt(15)
        para.font.color.rgb = WM_SLATE
        para.space_after = Pt(8)

    slide.notes_slide.notes_text_frame.text = _notes(notes)


def build_presentation(
    slug: str,
    fm: dict[str, str],
    body: str,
    *,
    out_path: Path | None = None,
    locale: str = "en",
) -> Path:
    copy = COPY.get(locale, COPY["en"])
    title = fm.get("title", slug.replace("-", " ").title())
    state = fm.get("state", "").strip()
    category = fm.get("category", "Car Accidents")
    excerpt = fm.get("excerpt", fm.get("description", ""))[:400]
    date = fm.get("date", "")
    author_id = fm.get("authorId", "scott-tischler")
    author_name = "Kathy Carr, CEO" if author_id == "kathy-carr" else "Scott Tischler, Co-Founder"

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    url = f"{SITE}/es/blog/{slug}" if locale == "es" else f"{SITE}/blog/{slug}"
    subtitle = f"{category}" + (f" · {state}" if state else "") + (f" · {date}" if date else "")

    _add_title_slide(
        prs,
        title,
        subtitle,
        f"Educational presentation for {title}. Full guide: {url}. WreckMatch LLC legal referral service, not a law firm.",
    )

    disclaimer_bullets = [
        copy["disclaimer_1"],
        copy["disclaimer_2"],
        copy["disclaimer_3"],
        copy["disclaimer_4"],
        f"Author: {author_name}.",
    ]
    _add_content_slide(prs, copy["read_first"], disclaimer_bullets, copy["read_first"])

    quick = find_quick_answer(body) or excerpt
    if quick:
        _add_content_slide(prs, copy["quick_answer"], [quick], quick)

    sections = extract_sections(body)
    content_slides = 0
    for sec in sections:
        if content_slides >= 8:
            break
        if re.search(r"why we published|note for families|por qué publicamos", sec.heading, re.I):
            _add_content_slide(prs, sec.heading, sec.bullets[:4], sec.notes)
            content_slides += 1
            continue
        if len(sec.bullets) >= 1:
            _add_content_slide(prs, sec.heading, sec.bullets, sec.notes)
            content_slides += 1

    pairs = faq_pairs(body)
    if pairs:
        faq_bullets: list[str] = []
        for q, a in pairs[:5]:
            faq_bullets.append(f"Q: {q}")
            faq_bullets.append(f"A: {a}")
        _add_content_slide(prs, copy["faq_heading"], faq_bullets, copy["faq_heading"])

    cta_bullets = [
        copy["cta_1"].format(site=SITE, phone=PHONE, url=url),
        copy["cta_2"].format(site=SITE, phone=PHONE, url=url),
        copy["cta_3"].format(site=SITE, phone=PHONE, url=url),
        copy["cta_4"].format(site=SITE, phone=PHONE, url=url),
    ]
    _add_content_slide(prs, copy["cta_heading"], cta_bullets, copy["cta_heading"])

    _add_content_slide(
        prs,
        copy["review_heading"],
        [
            "Judge Roy Waddell — Legal Advisor, WreckMatch LLC.",
            "Courtroom and procedural perspective only.",
            "Not case-specific legal advice.",
            f"Guía: {url}",
        ],
        copy["review_heading"],
    )

    prs.core_properties.title = title[:240]
    prs.core_properties.subject = f"WreckMatch educational guide — {category}"
    prs.core_properties.category = "Legal Education"
    prs.core_properties.keywords = ", ".join(
        filter(
            None,
            [
                "WreckMatch",
                "car accident",
                state,
                category,
                slug.replace("-", " "),
                "personal injury",
                "attorney matching",
            ],
        )
    )[:240]
    prs.core_properties.comments = (
        f"Gold-tier presentation summary for {url}. Educational only; not a law firm."
    )

    dest = out_path or (
        PRESENTATIONS_DIR / "es" / f"{slug}.pptx"
        if locale == "es"
        else PRESENTATIONS_DIR / f"{slug}.pptx"
    )
    dest.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dest))
    return dest


def score_presentation(slug: str, path: Path, fm: dict[str, str], body: str) -> PresentationReport:
    issues: list[str] = []
    score = 100

    if not path.exists() or path.stat().st_size < 8000:
        issues.append("missing_or_tiny_file")
        return PresentationReport(slug=slug, score=0, slide_count=0, issues=issues, path=path)

    try:
        prs = Presentation(str(path))
        slide_count = len(prs.slides)
    except Exception:
        issues.append("corrupt_pptx")
        return PresentationReport(slug=slug, score=0, slide_count=0, issues=issues, path=path)

    if slide_count < 10:
        issues.append("below_min_slides")
        score -= min(40, (10 - slide_count) * 8)

    notes_count = sum(
        1
        for s in prs.slides
        if (s.notes_slide.notes_text_frame.text or "").strip()
    )
    if notes_count < slide_count:
        issues.append("missing_speaker_notes")
        score -= 15

    props = prs.core_properties
    if not (props.title and props.subject and props.keywords):
        issues.append("weak_metadata")
        score -= 10

    combined = " ".join(
        (s.notes_slide.notes_text_frame.text or "") for s in prs.slides
    ).lower()
    for marker, label in [
        ("not a law firm", "disclaimer"),
        ("800+", "network"),
        ("855", "phone"),
        ("wreckmatch.com", "url"),
    ]:
        if marker not in combined:
            issues.append(f"missing_{label}")
            score -= 12

    if "judge roy waddell" not in combined:
        issues.append("missing_roy")
        score -= 5

    if faq_pairs(body) and not any(
        token in combined
        for token in (
            "frequently asked",
            "preguntas frecuentes",
            "faq",
            "q:",
            "p:",
        )
    ):
        issues.append("missing_faq_content")
        score -= 8

    score = max(0, min(100, score))
    return PresentationReport(
        slug=slug,
        score=score,
        slide_count=slide_count,
        issues=issues,
        path=path,
    )


def generate_for_post(path: Path, *, force: bool = False, locale: str = "en") -> PresentationReport:
    slug = path.stem
    text = path.read_text(encoding="utf-8")
    fm, body = parse_frontmatter(text)
    out = (
        PRESENTATIONS_DIR / "es" / f"{slug}.pptx"
        if locale == "es"
        else PRESENTATIONS_DIR / f"{slug}.pptx"
    )

    if out.exists() and not force:
        return score_presentation(slug, out, fm, body)

    build_presentation(slug, fm, body, out_path=out, locale=locale)
    return score_presentation(slug, out, fm, body)


def upsert_frontmatter_presentation(path: Path, presentation_url: str) -> bool:
    text = path.read_text(encoding="utf-8")
    if not text.startswith("---"):
        return False
    parts = text.split("---", 2)
    if len(parts) < 3:
        return False
    fm, body = parts[1], parts[2]
    line = f'presentationUrl: "{presentation_url}"\n'
    if re.search(r"^presentationUrl:", fm, re.M):
        new_fm = re.sub(
            r"^presentationUrl:\s*(?:>[-]?\s*\n(?:  .+\n)+|[^\n]+)",
            line.strip(),
            fm,
            count=1,
            flags=re.M,
        )
    else:
        new_fm = fm.rstrip() + "\n" + line
    if new_fm == fm:
        return False
    path.write_text(f"---{new_fm}---{body}", encoding="utf-8")
    return True
